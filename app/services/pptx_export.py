"""최종 기획서(Markdown) → PowerPoint(.pptx) 변환.

발표용 산출물. `# 제목`은 표지 슬라이드, 각 `## 섹션`은 내용 슬라이드가 되고,
마크다운 표(PESTEL 등)는 네이티브 PPTX 표로, 목록(-)·문단·`**굵게**`는 본문 텍스트로 렌더한다.
python-pptx만 사용(외부 변환기 불필요)하며, 한 슬라이드에 내용이 넘치면 "(계속)" 슬라이드로 분할한다.
"""
from __future__ import annotations

import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_DIR = Path("outputs")

_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_TABLE_ROW = re.compile(r"^\s*\|(.+)\|\s*$")
_TABLE_SEP = re.compile(r"^\s*\|[\s:|-]+\|\s*$")
_HEADING = re.compile(r"^\s*(#{1,6})\s+(.*)")
_BULLET = re.compile(r"^\s*[-*]\s+")

# 16:9 슬라이드 · 본문 영역(EMU는 Inches로 표현)
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_MARGIN_L = Inches(0.7)
_BODY_W = Inches(11.93)          # 13.333 - 좌우 여백 0.7
_BODY_TOP = Inches(1.55)         # 제목 아래
_BODY_BOTTOM = Inches(7.05)      # 하단 여백 전 한계선

# 팔레트 (index.html 톤과 맞춘 차분한 블루/그레이)
_ACCENT = RGBColor(0x2B, 0x59, 0xC3)
_INK = RGBColor(0x1F, 0x2A, 0x37)
_HEAD_BG = RGBColor(0x2B, 0x59, 0xC3)
_HEAD_FG = RGBColor(0xFF, 0xFF, 0xFF)
_ROW_ALT = RGBColor(0xEE, 0xF2, 0xFB)


def _slugify(name: str) -> str:
    slug = re.sub(r"[^0-9A-Za-z가-힣]+", "-", name).strip("-")
    return slug or "plan"


def _cells(line: str) -> list[str]:
    return [c.strip() for c in _TABLE_ROW.match(line).group(1).split("|")]


def _parse_blocks(markdown: str) -> list[dict]:
    """마크다운을 블록 리스트로 파싱: heading/table/bullet/paragraph."""
    lines = (markdown or "").split("\n")
    blocks: list[dict] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        # 표
        if _TABLE_ROW.match(line) and i + 1 < len(lines) and _TABLE_SEP.match(lines[i + 1]):
            header = _cells(line)
            i += 2
            rows = []
            while i < len(lines) and _TABLE_ROW.match(lines[i]) and not _TABLE_SEP.match(lines[i]):
                rows.append(_cells(lines[i]))
                i += 1
            blocks.append({"kind": "table", "header": header, "rows": rows})
            continue
        m = _HEADING.match(line)
        if m:
            blocks.append({"kind": "heading", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1
            continue
        if _BULLET.match(line):
            blocks.append({"kind": "bullet", "text": _BULLET.sub("", line).strip()})
            i += 1
            continue
        if line.strip():
            blocks.append({"kind": "paragraph", "text": line.strip()})
        i += 1
    return blocks


def _add_runs(paragraph, text: str, size: int, color: RGBColor, bold: bool = False) -> None:
    """`**굵게**`를 인식해 run을 나눠 추가한다."""
    def _run(t: str, strong: bool) -> None:
        r = paragraph.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold or strong

    pos = 0
    for mm in _BOLD_RE.finditer(text):
        if mm.start() > pos:
            _run(text[pos:mm.start()], False)
        _run(mm.group(1), True)
        pos = mm.end()
    if pos < len(text):
        _run(text[pos:], False)


def _est_text_height(blocks: list[dict]) -> Inches:
    """텍스트 블록 묶음의 대략 높이(줄바꿈 추정). 페이지네이션 판단용."""
    total = 0.0
    for b in blocks:
        chars = len(b.get("text", ""))
        lines = max(1, -(-chars // 46))          # 한글 기준 줄당 ~46자
        per = 0.30 if b["kind"] == "bullet" else 0.32
        if b["kind"] == "heading":
            per = 0.40
        total += lines * per + 0.06
    return Inches(total + 0.05)


def _est_table_height(block: dict) -> Inches:
    n = 1 + len(block["rows"])
    return Inches(n * 0.36 + 0.15)


class _Deck:
    """제목 슬라이드 + 내용 슬라이드를 쌓고, 본문이 넘치면 자동 분할한다."""

    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self.slide = None
        self.title = ""
        self.top = _BODY_TOP

    def _content_layout(self):
        # "Title Only" 레이아웃(제목만) — 본문 영역을 직접 배치한다.
        return self.prs.slide_layouts[5]

    def new_slide(self, title: str, continued: bool = False) -> None:
        self.slide = self.prs.slides.add_slide(self._content_layout())
        self.title = title
        shown = f"{title} (계속)" if continued else title
        if self.slide.shapes.title is not None:
            self.slide.shapes.title.text = shown
            tf = self.slide.shapes.title.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(26)
                    r.font.bold = True
                    r.font.color.rgb = _INK
        self.top = _BODY_TOP

    def _fits(self, height: Inches) -> bool:
        return self.top + height <= _BODY_BOTTOM

    def _ensure(self, height: Inches) -> None:
        if self.slide is None:
            self.new_slide(self.title or "내용")
        # 이미 뭔가 그렸고 이번 블록이 안 들어가면 (계속) 슬라이드로
        elif self.top > _BODY_TOP and not self._fits(height):
            self.new_slide(self.title, continued=True)

    def add_text(self, blocks: list[dict]) -> None:
        """텍스트 블록을 현재 슬라이드에 채우고, 남는 높이를 넘기면 (계속) 슬라이드로 분할한다."""
        if not blocks:
            return
        if self.slide is None:
            self.new_slide(self.title or "내용")
        i, n = 0, len(blocks)
        while i < n:
            # 현재 슬라이드에 들어갈 만큼 블록을 그리디하게 모은다(최소 1개는 반드시 배치)
            chunk, used = [], 0.0
            while i < n:
                bh = _est_text_height([blocks[i]]).inches
                if chunk and self.top.inches + used + bh > _BODY_BOTTOM.inches:
                    break
                chunk.append(blocks[i])
                used += bh
                i += 1
            self._render_text_chunk(chunk, Inches(used))
            if i < n:                                   # 남은 블록 → (계속) 슬라이드
                self.new_slide(self.title, continued=True)

    def _render_text_chunk(self, blocks: list[dict], height: Inches) -> None:
        box = self.slide.shapes.add_textbox(_MARGIN_L, self.top, _BODY_W, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        first = True
        for b in blocks:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(4)
            if b["kind"] == "bullet":
                r0 = p.add_run()
                r0.text = "• "
                r0.font.size = Pt(15)
                r0.font.color.rgb = _ACCENT
                _add_runs(p, b["text"], 15, _INK)
            elif b["kind"] == "heading":       # H3+ → 소제목
                _add_runs(p, b["text"], 17, _ACCENT, bold=True)
                p.space_before = Pt(6)
            else:
                _add_runs(p, b["text"], 15, _INK)
        self.top = Inches(self.top.inches + height.inches)

    def add_table(self, block: dict) -> None:
        header, rows = block["header"], block["rows"]
        ncols = len(header)
        if ncols == 0:
            return
        h = _est_table_height(block)
        self._ensure(h)
        nrows = 1 + len(rows)
        gfx = self.slide.shapes.add_table(nrows, ncols, _MARGIN_L, self.top, _BODY_W, h)
        table = gfx.table
        for j, htxt in enumerate(header):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _HEAD_BG
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            _add_runs(para, htxt, 11, _HEAD_FG, bold=True)
        for ri, row in enumerate(rows, start=1):
            for j in range(ncols):
                cell = table.cell(ri, j)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _ROW_ALT
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                para = cell.text_frame.paragraphs[0]
                _add_runs(para, row[j] if j < len(row) else "", 11, _INK)
        self.top = self.top + h + Inches(0.15)


def build_pptx(markdown: str, title: str = "") -> Presentation:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blocks = _parse_blocks(markdown)

    # 표지: 첫 H1(있으면) 또는 넘겨받은 title
    cover_title = title
    if blocks and blocks[0]["kind"] == "heading" and blocks[0]["level"] == 1:
        cover_title = blocks[0]["text"]
        blocks = blocks[1:]
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = cover_title or "서비스 기획서"
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "AI 서비스 기획 보조 Multi-Agent · 자동 생성 기획서"

    deck = _Deck(prs)
    text_buf: list[dict] = []

    def flush() -> None:
        deck.add_text(text_buf)
        text_buf.clear()

    for b in blocks:
        if b["kind"] == "heading" and b["level"] <= 2:
            flush()
            deck.new_slide(b["text"])
        elif b["kind"] == "table":
            flush()
            deck.add_table(b)
        else:
            text_buf.append(b)   # heading level>=3, bullet, paragraph
    flush()
    return prs


def pptx_bytes(markdown: str, title: str = "") -> bytes:
    buf = io.BytesIO()
    build_pptx(markdown, title).save(buf)
    return buf.getvalue()


def save_pptx(project_name: str, markdown: str) -> str:
    OUTPUT_DIR.mkdir(exist_ok=True)
    path = OUTPUT_DIR / f"{_slugify(project_name)}.pptx"
    build_pptx(markdown, title=project_name).save(str(path))
    return str(path)
