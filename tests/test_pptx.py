"""Markdown → PPTX 변환 및 /export/pptx 엔드포인트 테스트."""
from __future__ import annotations

import io

from fastapi.testclient import TestClient
from pptx import Presentation

from app.main import app
from app.services import pptx_export

MD = """# 테스트 기획서

## 프로젝트 개요
**핵심** 요약 문장.
- 첫 번째 근거
- 두 번째 근거

## PESTEL 분석
| 요인 | 주요 내용 | 기회 | 위협 | 대응 |
|---|---|---|---|---|
| Political | 규제 | 지원 | 불확실 | 모니터링 |
| Economic | 성장 | 수요 | 비용 | 효율화 |

## 참고자료
- 출처 A — https://a.io
- https://b.io
"""


def _texts(slide) -> str:
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
    return "\n".join(out)


def test_build_pptx_structure():
    prs = pptx_export.build_pptx(MD)
    slides = list(prs.slides)
    # 표지 1 + 섹션 3(개요/PESTEL/참고자료) 이상
    assert len(slides) >= 4
    # 표지 제목은 H1
    assert "테스트 기획서" in _texts(slides[0])
    # 어딘가에 PESTEL 제목 슬라이드
    all_text = "\n".join(_texts(s) for s in slides)
    assert "PESTEL 분석" in all_text and "프로젝트 개요" in all_text
    # 표가 1개 렌더됨 (5열 x 3행)
    tables = [sh.table for s in slides for sh in s.shapes if sh.has_table]
    assert len(tables) == 1
    t = tables[0]
    assert len(t.columns) == 5 and len(t.rows) == 3
    assert t.cell(0, 0).text == "요인"
    assert t.cell(1, 0).text == "Political"


def test_16_9_dimensions():
    prs = pptx_export.build_pptx(MD)
    assert prs.slide_width == pptx_export._SLIDE_W
    assert prs.slide_height == pptx_export._SLIDE_H


def test_long_section_paginates():
    bullets = "\n".join(f"- 근거 항목 {i}: " + ("내용 " * 12) for i in range(40))
    md = "# 긴 문서\n\n## 아주 긴 섹션\n" + bullets
    prs = pptx_export.build_pptx(md)
    cont = [s for s in prs.slides if "(계속)" in _texts(s)]
    assert cont, "긴 섹션은 (계속) 슬라이드로 분할되어야 함"


def test_pptx_bytes_is_valid_zip_openable():
    data = pptx_export.pptx_bytes(MD)
    assert data[:2] == b"PK"                       # pptx = zip 컨테이너
    Presentation(io.BytesIO(data))                  # 다시 열려야 유효


def test_export_pptx_endpoint():
    c = TestClient(app)
    r = c.post("/export/pptx", json={"project_name": "테스트", "markdown": MD})
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    assert r.content[:2] == b"PK"


def test_table_followed_by_body_text_in_same_section():
    """표 **다음에 본문이 오는** 문서가 변환된다 (2026-07-28 실제 500 회귀).

    `Inches()` 는 int 서브클래스(Emu)이고 파이썬 `+` 는 서브클래스를 유지하지 않으므로,
    `add_table` 이 `self.top` 을 순수 int 로 만들면 뒤따르는 `_render_text_chunk` 의
    `self.top.inches` 가 AttributeError 로 터진다.

    기존 MD 는 표 바로 뒤가 새 H2(=새 슬라이드로 top 초기화)여서 이 조합을 비껴갔다.
    실 LLM KOSENA 문서(표 뒤에 설명 문단이 이어짐)에서 `/export/pptx` 가 500 이 됐다.
    """
    md = (
        "# 문서\n\n## 표가 있는 섹션\n"
        "| 항목 | 값 |\n|---|---|\n| A | 1 |\n| B | 2 |\n\n"
        "표 아래에 이어지는 설명 문단이다.\n"
        "- 표 뒤 불릿 하나\n"
    )
    prs = pptx_export.build_pptx(md)                 # 여기서 터지면 회귀
    body = "\n".join(_texts(s) for s in prs.slides)
    assert "표 아래에 이어지는 설명 문단이다." in body
    assert any(sh.has_table for s in prs.slides for sh in s.shapes)

    c = TestClient(app)
    r = c.post("/export/pptx", json={"project_name": "표뒤본문", "markdown": md})
    assert r.status_code == 200
